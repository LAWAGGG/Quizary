"""Generate draf form (sections + soal + settings) via Gemini.

Alur: prompt + file referensi (teks) -> Gemini (JSON mode) -> draf yang
sudah disanitasi. Draf TIDAK langsung jadi form — creator mereview lalu
POST /api/ai/accept yang memvalidasi ulang memakai skema existing.
"""
import io
import json
import logging

import httpx

from app.config import GEMINI_API_KEY, GEMINI_FALLBACK_MODEL, GEMINI_MODEL
from app.schemas.question import QuestionCreate

logger = logging.getLogger("quizary.ai")

AI_DAILY_LIMIT = 3

ALLOWED_REF_EXT = {".docx", ".pdf", ".pptx"}
MAX_REF_FILES = 5
MAX_REF_FILE_BYTES = 5 * 1024 * 1024  # 5 MB per file
MAX_REF_TOTAL_CHARS = 30_000

MAX_SECTIONS = 10
MAX_QUESTIONS = 50
MAX_OPTIONS = 10

GEMINI_TIMEOUT = 120.0

QUESTION_TYPES = (
    "multiple_choice", "checkbox", "dropdown", "short_answer", "essay",
    "password", "date", "time", "datetime", "file_upload",
)
OPTION_TYPES = ("multiple_choice", "checkbox", "dropdown")
# quiz multiple_choice wajib tepat 1 kunci — dicek ulang saat accept.


class AiNotConfigured(Exception):
    pass


class AiFailed(Exception):
    pass


def extract_ref_text(filename: str, raw: bytes) -> str:
    """Ambil teks dari file referensi (docx/pdf/pptx)."""
    name = (filename or "").lower()
    if name.endswith(".docx"):
        from docx import Document  # type: ignore

        doc = Document(io.BytesIO(raw))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n".join(parts)
    if name.endswith(".pdf"):
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(io.BytesIO(raw))
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    if name.endswith(".pptx"):
        from pptx import Presentation  # type: ignore

        prs = Presentation(io.BytesIO(raw))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            parts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n".join(parts)
    raise AiFailed(f"Tipe file tidak didukung ({filename}). Pakai docx, pdf, atau pptx.")


SYSTEM_INSTRUCTION = """Kamu penyusun form/kuis berbahasa Indonesia. Jawab HANYA dengan SATU objek JSON valid, tanpa markdown, tanpa penjelasan.

Bentuk:
{"sections": [{"title": "nama section", "questions": [
  {"type": "salah satu: %s", "question_text": "teks soal",
   "is_required": true, "points": 1,
   "options": [{"option_text": "teks opsi", "is_correct": true}],
   "password_keyword": null}
]}], "settings": {"shuffle_questions": false, "shuffle_options": false, "timer_minutes": null}}

Aturan:
- options HANYA untuk multiple_choice/checkbox/dropdown (2-4 opsi); tipe lain: options [] dan password_keyword null.
- password_keyword HANYA untuk type password (isi kata sandinya), selain itu null.
- Quiz: multiple_choice WAJIB tepat 1 option is_correct=true; checkbox boleh >1; timer_minutes WAJIB angka 1-1440.
- Bukan quiz: timer_minutes null, is_correct semua false.
- Maksimal 10 sections, total maksimal 30 soal. question_text ringkas, tanpa HTML.
""" % (", ".join(QUESTION_TYPES))


def build_user_text(title: str, description: str | None, form_type: str, prompt: str, refs: list[tuple[str, str]]) -> str:
    parts = [
        f"Jenis: {'KUIS (ada nilai & kunci jawaban)' if form_type == 'quiz' else 'FORMULIR/pendataan (tanpa nilai)'}",
        f"Judul: {title}",
    ]
    if description:
        parts.append(f"Deskripsi: {description}")
    parts.append(f"Permintaan creator:\n{prompt}")
    for fname, text in refs:
        parts.append(f"--- Isi file referensi {fname} ---\n{text}")
    return "\n\n".join(parts)


def _gemini_models() -> list[str]:
    """Model utama + cadangan (dedupe). Kosong = tanpa fallback."""
    models = [m.strip() for m in (GEMINI_MODEL, GEMINI_FALLBACK_MODEL) if m and m.strip()]
    return list(dict.fromkeys(models)) or ["gemini-3.6-flash"]


def _parse_gemini_text(data: dict) -> dict:
    if data.get("promptFeedback", {}).get("blockReason"):
        raise AiFailed("Prompt ditolak filter keamanan AI. Coba ubah kata-katanya.")
    cands = data.get("candidates") or []
    text = (cands[0].get("content", {}).get("parts") or [{}])[0].get("text", "") if cands else ""
    try:
        parsed = json.loads(text)
    except (ValueError, TypeError):
        raise AiFailed("AI gagal menyusun draf. Coba generate ulang.")
    if not isinstance(parsed, dict):
        raise AiFailed("AI gagal menyusun draf. Coba generate ulang.")
    return parsed


def call_gemini(user_text: str) -> tuple[dict, str]:
    """Panggil Gemini JSON mode. Balik (draf, model_terpakai).

    Key dikirim via header (tak muncul di URL/log). Tiap model dicoba 2x
    untuk error transient (429/5xx/network/JSON rusak); gagal semua di model
    utama -> lanjut ke fallback. 401/403 (key salah) dan 400 langsung gagal
    tanpa buang kuota coba — fallback pakai key yang sama.
    """
    if not GEMINI_API_KEY:
        raise AiNotConfigured("Fitur AI belum dikonfigurasi server.")
    payload = {
        "systemInstruction": {"parts": [{"text": SYSTEM_INSTRUCTION}]},
        "contents": [{"parts": [{"text": user_text}]}],
        "generationConfig": {"responseMimeType": "application/json", "temperature": 0.7, "maxOutputTokens": 8192},
    }
    headers = {"x-goog-api-key": GEMINI_API_KEY}
    last_err: Exception | None = None
    for model in _gemini_models():
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
        for _ in range(2):
            try:
                with httpx.Client(timeout=GEMINI_TIMEOUT) as client:
                    resp = client.post(url, headers=headers, json=payload)
            except httpx.HTTPError as e:
                logger.warning("gemini %s: network error %s", model, type(e).__name__)
                last_err = e
                continue
            if resp.status_code == 200:
                try:
                    return _parse_gemini_text(resp.json()), model
                except AiFailed as e:
                    logger.warning("gemini %s: draf tak valid (%s)", model, e)
                    last_err = e
                    continue
                except ValueError as e:
                    logger.warning("gemini %s: respons bukan JSON", model)
                    last_err = e
                    continue
            if resp.status_code in (401, 403):
                logger.error("gemini %s: key ditolak (%s)", model, resp.status_code)
                raise AiFailed("API key AI ditolak. Hubungi admin.")
            if resp.status_code == 400:
                logger.warning("gemini %s: 400 %.120s", model, resp.text)
                raise AiFailed("AI menolak permintaan. Coba ubah prompt lalu generate ulang.")
            if resp.status_code == 404:
                # ID model pensiun/diganti Google (kasus 2.x) — bukan "sibuk".
                logger.error("gemini %s: 404 model tak tersedia", model)
                last_err = AiFailed(f"Model AI {model} tidak tersedia. Hubungi admin.")
                break
            # 429 / 5xx -> coba lagi / fallback.
            logger.warning("gemini %s: sibuk (%s)", model, resp.status_code)
            last_err = AiFailed(f"AI sibuk ({resp.status_code}). Coba lagi sebentar lagi.")
    if isinstance(last_err, AiFailed):
        raise last_err
    raise AiFailed("AI tidak merespons. Periksa koneksi lalu coba lagi.")


def _coerce_question(raw: dict) -> dict | None:
    """Bersihkan 1 soal AI -> dict valid QuestionCreate, atau None bila sampah."""
    if not isinstance(raw, dict):
        return None
    q_type = raw.get("type") if raw.get("type") in QUESTION_TYPES else None
    text = str(raw.get("question_text") or "").strip()
    if not q_type or not text:
        return None
    text = text[:5000]
    opts: list[dict] = []
    if q_type in OPTION_TYPES:
        for o in (raw.get("options") or [])[:MAX_OPTIONS]:
            if not isinstance(o, dict):
                continue
            t = str(o.get("option_text") or "").strip()
            if t:
                opts.append({"option_text": t[:2000], "is_correct": bool(o.get("is_correct"))})
        if not opts:
            return None
    try:
        points = int(raw.get("points", 1))
    except (TypeError, ValueError):
        points = 1
    kw = str(raw.get("password_keyword") or "").strip() or None
    try:
        q = QuestionCreate(
            type=q_type,
            question_text=text,
            points=max(0, min(999, points)),
            is_required=bool(raw.get("is_required", True)),
            password_keyword=kw if q_type == "password" else None,
            options=opts,
        )
    except Exception:
        return None
    return q.model_dump()


def sanitize_draft(raw: dict, form_type: str) -> dict:
    """Bersihkan output AI -> draf valid. Raise AiFailed bila tak ada soal layak."""
    sections: list[dict] = []
    total = 0
    for s in (raw.get("sections") or [])[:MAX_SECTIONS]:
        title = str((s or {}).get("title") or "").strip()[:150] or "Bagian"
        questions: list[dict] = []
        for q in ((s or {}).get("questions") or []):
            if total >= MAX_QUESTIONS:
                break
            clean = _coerce_question(q)
            if clean:
                questions.append(clean)
                total += 1
        if questions:
            sections.append({"title": title, "questions": questions})
    if not sections:
        raise AiFailed("AI tidak menghasilkan soal yang valid. Coba perjelas prompt lalu generate ulang.")

    settings = raw.get("settings") or {}
    try:
        timer = settings.get("timer_minutes")
        timer = int(timer) if timer is not None else None
        timer = timer if timer is not None and 1 <= timer <= 1440 else None
    except (TypeError, ValueError):
        timer = None
    sub = settings.get("submission_limit")
    draft = {
        "sections": sections,
        "settings": {
            "shuffle_questions": bool(settings.get("shuffle_questions", False)),
            "shuffle_options": bool(settings.get("shuffle_options", False)),
            "timer_minutes": timer,
            "require_login": bool(settings.get("require_login", False)),
            "submission_limit": sub if sub in ("unlimited", "once") else "unlimited",
        },
    }
    if form_type == "quiz" and timer is None:
        # Jangan diam-diam tanpa timer (publish quiz wajib timer) — creator regenerate.
        raise AiFailed("AI tidak menyertakan timer untuk kuis. Coba generate ulang.")
    return draft
